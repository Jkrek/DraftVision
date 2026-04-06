"""
DraftVision trifold expo poster — 48" wide × 36" tall.
Designed for a physical trifold board with 12" left panel, 24" center panel,
and 12" right panel, printed as a single sheet.
"""

from __future__ import annotations

import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
KEYNOTE_PATH = ROOT / "DraftVision_Poster_v4.key"
OUT_PATH = ROOT / "DraftVision_Poster_trifold_36x48.pptx"

prs = Presentation()
prs.slide_width = Inches(48)
prs.slide_height = Inches(36)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BG = RGBColor(0x04, 0x08, 0x14)
HEADER = RGBColor(0x02, 0x05, 0x10)
PANEL = RGBColor(0x0C, 0x17, 0x2A)
CARD = RGBColor(0x12, 0x22, 0x3A)
CARD2 = RGBColor(0x0A, 0x18, 0x2E)
BLUE = RGBColor(0x38, 0x8E, 0xFF)
LBLUE = RGBColor(0x7A, 0xB8, 0xFF)
GREEN = RGBColor(0x0F, 0xE0, 0x7A)
GOLD = RGBColor(0xFF, 0xB3, 0x00)
RED = RGBColor(0xFF, 0x45, 0x45)
PURPLE = RGBColor(0xB0, 0x8A, 0xFF)
CYAN = RGBColor(0x00, 0xE5, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xD4, 0xE6, 0xF8)
DIM = RGBColor(0x78, 0x96, 0xB6)
VDIM = RGBColor(0x2E, 0x44, 0x5C)
UC_RED = RGBColor(0xE0, 0x00, 0x00)
UC_GLD = RGBColor(0xF5, 0xBF, 0x00)

LEFT_X = 0.0
CENTER_X = 12.0
RIGHT_X = 36.0
LEFT_W = 12.0
CENTER_W = 24.0
RIGHT_W = 12.0


def rect(l, t, w, h, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(l),
        Inches(t),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    return shape


def textbox(l, t, w, h, text, pt, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def bullet_box(l, t, w, items, accent, title=None, body_size=24, title_size=30, height=3.6):
    rect(l, t, w, height, CARD)
    rect(l, t, w, 0.10, accent)
    if title:
        textbox(l + 0.22, t + 0.14, w - 0.44, 0.42, title, title_size, accent, bold=True)
        top = t + 0.68
    else:
        top = t + 0.22

    tb = slide.shapes.add_textbox(Inches(l + 0.22), Inches(top), Inches(w - 0.44), Inches(height - (top - t) - 0.18))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.size = Pt(body_size)
        p.font.color.rgb = OFFWHT
        p.space_after = Pt(5)
    return tb


def section_bar(l, t, w, label, accent):
    rect(l, t, w, 0.82, CARD2)
    rect(l, t, 0.18, 0.82, accent)
    rect(l + 0.18, t, w - 0.18, 0.06, accent)
    textbox(l + 0.34, t + 0.10, w - 0.48, 0.56, label, 34, accent, bold=True)


def stat_card(l, t, w, h, value, label, accent):
    rect(l, t, w, h, CARD)
    rect(l, t, w, 0.11, accent)
    textbox(l, t + 0.18, w, 0.78, value, 52, accent, bold=True, align=PP_ALIGN.CENTER)
    textbox(l, t + 1.06, w, 0.40, label, 20, DIM, align=PP_ALIGN.CENTER)


def bar(l, t, w, label, value, accent, max_value=40):
    textbox(l, t - 0.02, w * 0.52, 0.36, label, 21, OFFWHT)
    rect(l + w * 0.52, t + 0.09, w * 0.34, 0.16, VDIM)
    fill = (w * 0.34) * min(value / max_value, 1.0)
    if fill > 0:
        rect(l + w * 0.52, t + 0.09, fill, 0.16, accent)
    textbox(l + w * 0.88, t - 0.06, w * 0.12, 0.32, f"{value}%", 18, accent, bold=True, align=PP_ALIGN.RIGHT)


def add_image(path: Path, l, t, w, h):
    return slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))


def extract_assets() -> dict[str, Path]:
    assets: dict[str, Path] = {}
    if not KEYNOTE_PATH.exists():
        return assets

    temp_dir = Path(tempfile.mkdtemp(prefix="draftvision_trifold_"))
    targets = {
        "headshot": "Facetune_16-02-2026-15-33-07-23.jpeg",
        "football": "Screenshot 2026-04-06 at 10.49-54.png",
        "website": "Screenshot 2026-04-06 at 10.49.48",
    }

    with zipfile.ZipFile(KEYNOTE_PATH) as zf:
        for info in zf.infolist():
            name = info.filename
            for key, target in targets.items():
                if target in name:
                    suffix = Path(name).suffix or ".png"
                    out = temp_dir / f"{key}{suffix}"
                    out.write_bytes(zf.read(info))
                    assets[key] = out
    return assets


assets = extract_assets()

# Background and fold guides
rect(0, 0, 48, 36, BG)
for y in range(0, 37, 4):
    rect(0, y, 48, 0.025, RGBColor(0x0A, 0x16, 0x28))

rect(11.94, 0, 0.12, 36, RGBColor(0x10, 0x20, 0x38))
rect(35.94, 0, 0.12, 36, RGBColor(0x10, 0x20, 0x38))
textbox(11.2, 34.9, 1.6, 0.3, "fold", 14, VDIM, align=PP_ALIGN.CENTER)
textbox(35.2, 34.9, 1.6, 0.3, "fold", 14, VDIM, align=PP_ALIGN.CENTER)

# Header
rect(0, 0, 48, 4.4, HEADER)
rect(0, 4.08, 48, 0.10, BLUE)
rect(0, 4.18, 48, 0.05, GREEN)
textbox(19.3, -0.55, 12, 5.4, "D", 360, RGBColor(0x06, 0x10, 0x22), bold=True, align=PP_ALIGN.CENTER)

if "football" in assets:
    add_image(assets["football"], 0.35, 0.35, 11.2, 3.35)
    rect(0.35, 0.35, 11.2, 3.35, RGBColor(0x02, 0x05, 0x10), line=RGBColor(0x02, 0x05, 0x10))
    slide.shapes._spTree.remove(slide.shapes[-1]._element)
    # left photo fade blocks
    rect(0.35, 0.35, 2.2, 3.35, RGBColor(0x02, 0x05, 0x10))

# central header card for stronger readability and cleaner front-facing design
rect(12.05, 0.28, 24.00, 3.18, CARD2)
rect(12.05, 0.28, 24.00, 0.10, BLUE)
rect(12.05, 3.36, 24.00, 0.10, GREEN)

# title shadow + split color title treatment
textbox(14.08, 0.40, 8.9, 1.15, "Draft", 90, RGBColor(0x0B, 0x1A, 0x2E), bold=True, align=PP_ALIGN.RIGHT)
textbox(22.85, 0.40, 11.2, 1.15, "Vision", 90, RGBColor(0x0B, 0x1A, 0x2E), bold=True, align=PP_ALIGN.LEFT)
textbox(14.0, 0.32, 8.9, 1.15, "Draft", 90, WHITE, bold=True, align=PP_ALIGN.RIGHT)
textbox(22.77, 0.32, 11.2, 1.15, "Vision", 90, BLUE, bold=True, align=PP_ALIGN.LEFT)

rect(16.25, 1.56, 15.55, 0.52, CARD)
rect(16.25, 1.56, 0.14, 0.52, GOLD)
textbox(16.55, 1.66, 14.95, 0.32, "Machine Learning–Powered NFL Prospect Success Prediction", 22, OFFWHT, bold=True, align=PP_ALIGN.CENTER)

rect(14.15, 2.24, 9.95, 0.52, CARD)
rect(14.15, 2.24, 0.12, 0.52, CYAN)
textbox(14.40, 2.34, 9.45, 0.30, "Jared Krekeler  ·  University of Cincinnati", 17, OFFWHT, bold=True, align=PP_ALIGN.CENTER)

rect(24.50, 2.24, 9.40, 0.52, CARD)
rect(24.50, 2.24, 0.12, 0.52, GREEN)
textbox(24.78, 2.34, 8.85, 0.30, "Advisor: Hrishi Vinayak Bhide", 17, OFFWHT, bold=True, align=PP_ALIGN.CENTER)

textbox(15.10, 2.92, 17.8, 0.28, "Spring 2025  ·  CS Senior Design Expo", 16, DIM, italic=True, align=PP_ALIGN.CENTER)
rect(18.1, 3.52, 11.8, 0.08, BLUE)
rect(19.6, 3.64, 8.8, 0.05, GOLD)

if "headshot" in assets:
    add_image(assets["headshot"], 44.15, 0.42, 3.10, 3.10)
else:
    rect(44.15, 0.42, 3.10, 3.10, CARD)
rect(40.2, 0.42, 3.25, 1.42, UC_RED)
textbox(40.2, 0.48, 3.25, 0.60, "UC", 42, WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(40.2, 1.00, 3.25, 0.42, "CEAS", 20, WHITE, align=PP_ALIGN.CENTER)
rect(39.95, 2.02, 3.55, 1.02, CARD)
rect(39.95, 2.02, 3.55, 0.08, GREEN)
textbox(40.05, 2.16, 3.35, 0.28, "Trifold Layout", 15, GREEN, bold=True, align=PP_ALIGN.CENTER)
textbox(40.00, 2.45, 3.45, 0.42, "36\" × 48\" overall\n12\" | 24\" | 12\" panels", 14, OFFWHT, align=PP_ALIGN.CENTER)

# Panel backgrounds
body_top = 4.75
body_h = 27.0
rect(0.35, body_top, 11.2, body_h, PANEL)
rect(12.4, body_top, 23.2, body_h, PANEL)
rect(36.45, body_top, 11.2, body_h, PANEL)

# Left panel
section_bar(0.55, 4.95, 10.8, "PROBLEM & PIPELINE", BLUE)
textbox(0.82, 5.95, 10.2, 1.95, "NFL draft decisions are expensive, but prospect evaluation is still largely subjective. DraftVision adds a data-driven probability of NFL success to support scouting decisions.", 22, OFFWHT)

rect(0.82, 8.05, 9.65, 1.35, CARD)
rect(0.82, 8.05, 0.14, 1.35, GOLD)
textbox(1.18, 8.28, 9.0, 0.70, '"Only ~31% of drafted players become long-term NFL starters."', 22, GOLD, italic=True)

bullet_box(
    0.82,
    9.75,
    9.65,
    [
        "Prospect name entered in web app",
        "Live ESPN college roster + stat lookup",
        "15-feature vector assembled",
        "XGBoost model predicts NFL success probability",
        "Result returned with draft projection",
    ],
    CYAN,
    title="Pipeline",
    body_size=22,
    title_size=30,
    height=4.8,
)

bullet_box(
    0.82,
    14.95,
    9.65,
    [
        "Draft round",
        "Combine athleticism score",
        "College tier",
        "Production score",
        "Passing/rushing volume stats",
        "Position flags",
    ],
    PURPLE,
    title="Key Inputs",
    body_size=22,
    title_size=30,
    height=4.7,
)

bullet_box(
    0.82,
    20.05,
    9.65,
    [
        "No public labeled dataset",
        "Strong class imbalance",
        "Different positions produce very different stats",
        "Needed live data fast enough for a real web app",
    ],
    RED,
    title="Main Challenges",
    body_size=21,
    title_size=30,
    height=4.9,
)

bullet_box(
    0.82,
    25.35,
    9.65,
    [
        "Rounds 3–7 value scouting",
        "Fast fan/analyst exploration",
        "Repeatable, explainable predictions",
    ],
    GREEN,
    title="Why It Matters",
    body_size=21,
    title_size=30,
    height=3.7,
)

# Center panel
section_bar(12.75, 4.95, 22.5, "RESULTS & APPLICATION", GREEN)
stat_card(13.05, 6.05, 6.9, 1.72, "86.8%", "Accuracy", GREEN)
stat_card(20.55, 6.05, 6.9, 1.72, "0.928", "AUC", BLUE)
stat_card(28.05, 6.05, 6.9, 1.72, "0.85", "Precision", GOLD)

rect(13.05, 8.15, 10.8, 5.10, CARD)
rect(13.05, 8.15, 10.8, 0.10, GREEN)
textbox(13.35, 8.35, 10.1, 0.48, "Feature Importance", 30, GREEN, bold=True)
for idx, (label, value, accent) in enumerate([
    ("Draft Round", 34.3, GOLD),
    ("Combine Athleticism", 6.8, BLUE),
    ("Production Score", 5.5, GREEN),
    ("College Tier", 5.5, GREEN),
    ("Passing Yards", 5.4, DIM),
    ("Rushing Yards", 5.3, DIM),
]):
    bar(13.35, 9.00 + idx * 0.54, 9.4, label, value, accent)

rect(24.15, 8.15, 10.80, 5.10, CARD)
rect(24.15, 8.15, 10.80, 0.10, BLUE)
textbox(24.45, 8.35, 10.1, 0.48, "Prediction Output", 30, BLUE, bold=True)
for i, (name, val, color) in enumerate([
    ("NFL Success", "68%  →  Round 1", GREEN),
    ("Marginal", "26%  →  Round 4", GOLD),
    ("Low Probability", "4%   →  Undrafted", RED),
]):
    y = 9.02 + i * 1.23
    rect(24.45, y, 10.2, 0.98, CARD2)
    rect(24.45, y, 0.12, 0.98, color)
    textbox(24.75, y + 0.12, 4.2, 0.30, name, 21, DIM)
    textbox(24.75, y + 0.44, 8.8, 0.34, val, 23, color, bold=True)

rect(13.05, 13.65, 10.8, 7.15, CARD)
rect(13.05, 13.65, 10.8, 0.10, PURPLE)
textbox(13.35, 13.88, 10.0, 0.48, "Draft Round Success Rates", 30, PURPLE, bold=True)
for idx, (label, val, color) in enumerate([
    ("Round 1", 68, GREEN),
    ("Round 2", 52, LBLUE),
    ("Round 3", 38, BLUE),
    ("Round 4", 26, PURPLE),
    ("Round 5", 18, DIM),
    ("Round 6", 12, DIM),
    ("Round 7", 8, VDIM),
    ("Undrafted", 4, VDIM),
]):
    y = 14.55 + idx * 0.72
    rect(13.35, y, 8.0, 0.24, CARD2)
    rect(13.35, y, 8.0 * val / 100, 0.24, color)
    textbox(13.50, y - 0.08, 2.2, 0.34, label, 18, WHITE, bold=True)
    textbox(21.55, y - 0.08, 1.6, 0.34, f"{val}%", 18, color, bold=True, align=PP_ALIGN.RIGHT)

rect(24.15, 13.65, 10.80, 7.15, CARD)
rect(24.15, 13.65, 10.80, 0.10, CYAN)
textbox(24.45, 13.88, 10.0, 0.48, "Web Platform", 30, CYAN, bold=True)
textbox(24.45, 14.42, 9.95, 1.35, "Users enter any prospect's name and get a real-time AI scouting report with live stats, combine profile, success probability, and a draft projection.", 22, OFFWHT)
if "website" in assets:
    add_image(assets["website"], 24.52, 15.92, 9.95, 3.78)
else:
    rect(24.52, 15.92, 9.95, 3.78, CARD2)
    textbox(24.52, 17.10, 9.95, 0.6, "Website screenshot unavailable", 22, DIM, align=PP_ALIGN.CENTER)
textbox(24.55, 19.98, 9.9, 0.36, "Stack: Python · Flask · React · XGBoost · PostgreSQL · Railway", 17, DIM, italic=True, align=PP_ALIGN.CENTER)

rect(13.05, 21.20, 21.9, 8.20, CARD)
rect(13.05, 21.20, 21.9, 0.10, GOLD)
textbox(13.45, 21.44, 21.1, 0.52, "Takeaways", 32, GOLD, bold=True)
left_takeaway = slide.shapes.add_textbox(Inches(13.45), Inches(22.10), Inches(10.2), Inches(6.6))
ltf = left_takeaway.text_frame
ltf.word_wrap = True
ltf.clear()
for idx, item in enumerate([
    "XGBoost clearly outperformed baseline models.",
    "Draft round is the strongest predictor by a wide margin.",
    "The project works as a live web product, not just a notebook model.",
    "Future work: real historical NFL outcome labels and position-specific models.",
]):
    p = ltf.paragraphs[0] if idx == 0 else ltf.add_paragraph()
    p.text = item
    p.bullet = True
    p.font.size = Pt(23)
    p.font.color.rgb = OFFWHT
    p.space_after = Pt(8)

right_takeaway = slide.shapes.add_textbox(Inches(24.6), Inches(22.08), Inches(9.8), Inches(6.6))
rtf = right_takeaway.text_frame
rtf.word_wrap = True
rtf.clear()
for idx, item in enumerate([
    "Useful for scouting discussions in the middle and late rounds.",
    "Gives fans and analysts a simple explanation of prospect risk.",
    "Auto-updating ESPN inputs reduce manual maintenance.",
    "Designed larger for readable expo printing.",
]):
    p = rtf.paragraphs[0] if idx == 0 else rtf.add_paragraph()
    p.text = item
    p.bullet = True
    p.font.size = Pt(23)
    p.font.color.rgb = OFFWHT
    p.space_after = Pt(8)

# Right panel
section_bar(36.65, 4.95, 10.8, "METHODS", PURPLE)
bullet_box(
    36.95,
    5.95,
    9.65,
    [
        "ESPN CFB roster API",
        "ESPN season stats API",
        "4,000 synthetic labeled samples",
        "15 engineered features",
        "XGBoost classifier",
    ],
    PURPLE,
    title="Data + Model",
    body_size=22,
    title_size=30,
    height=4.8,
)

bullet_box(
    36.95,
    11.10,
    9.65,
    [
        "Logistic Regression: 72.0%",
        "Random Forest: 74.2%",
        "XGBoost: 86.8%",
    ],
    GREEN,
    title="Model Comparison",
    body_size=23,
    title_size=30,
    height=3.6,
)

bullet_box(
    36.95,
    15.05,
    9.65,
    [
        "Real outcome labels from historical NFL careers",
        "Position-specific models",
        "Tracking data such as separation and speed",
        "More explainability for each prediction",
    ],
    GOLD,
    title="Future Work",
    body_size=22,
    title_size=30,
    height=4.8,
)

rect(36.95, 20.25, 9.65, 4.25, CARD)
rect(36.95, 20.25, 9.65, 0.10, RED)
textbox(37.25, 20.48, 9.0, 0.46, "Expo Notes", 30, RED, bold=True)
textbox(37.25, 21.05, 8.95, 2.35, "This file is sized for a 36\" × 48\" trifold board printed as one full sheet. Center panel is 24\" wide. Side panels are 12\" wide each. Cut on the fold lines before mounting.", 22, OFFWHT)
textbox(37.25, 23.45, 8.95, 0.34, "Keep critical text away from the fold edges.", 18, DIM, italic=True)

rect(36.95, 24.90, 9.65, 4.15, CARD)
rect(36.95, 24.90, 9.65, 0.10, BLUE)
textbox(37.25, 25.12, 9.0, 0.44, "Contact", 30, BLUE, bold=True)
textbox(37.25, 25.72, 8.95, 2.10, "Jared Krekeler\nB.S. Computer Science\nUniversity of Cincinnati\nkrekeljm@mail.uc.edu", 24, OFFWHT)
textbox(37.25, 28.25, 8.95, 0.30, "Advisor: Hrishi Vinayak Bhide", 18, GREEN, bold=True)

# Footer
rect(0, 32.3, 48, 3.7, HEADER)
rect(0, 32.3, 48, 0.10, BLUE)
rect(0, 32.4, 48, 0.05, GREEN)
textbox(0.8, 32.78, 14.0, 0.38, "DraftVision  ·  CS Senior Design Expo", 19, WHITE, bold=True)
textbox(16.0, 32.78, 16.0, 0.38, "University of Cincinnati  ·  Spring 2025", 19, DIM, align=PP_ALIGN.CENTER)
textbox(33.4, 32.78, 13.6, 0.38, "Designed for 12\" | 24\" | 12\" trifold mounting", 18, GREEN, align=PP_ALIGN.RIGHT)
textbox(1.0, 33.35, 46.0, 0.44, "Print as one 36\" × 48\" sheet. Then cut and mount to the trifold board panels manually.", 21, OFFWHT, align=PP_ALIGN.CENTER)
rect(0, 35.82, 48, 0.10, UC_RED)
rect(0, 35.92, 48, 0.08, UC_GLD)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
