"""
DraftVision Expo Poster — 42" wide × 48" tall
v6: drastically larger text, trimmed content, and assets pulled from v4 Keynote.
"""

from __future__ import annotations

import os
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
KEYNOTE_PATH = ROOT / "DraftVision_Poster_v4.key"
OUT_PATH = ROOT / "DraftVision_Poster_v6.pptx"

prs = Presentation()
prs.slide_width = Inches(42)
prs.slide_height = Inches(48)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── Palette ──────────────────────────────────────────────────────────────────
BG = RGBColor(0x04, 0x08, 0x14)
PANEL = RGBColor(0x0C, 0x17, 0x2A)
CARD = RGBColor(0x12, 0x22, 0x3A)
CARD2 = RGBColor(0x0A, 0x18, 0x2E)
BLUE = RGBColor(0x38, 0x8E, 0xFF)
LBLUE = RGBColor(0x7A, 0xB8, 0xFF)
GREEN = RGBColor(0x0F, 0xE0, 0x7A)
GOLD = RGBColor(0xFF, 0xB3, 0x00)
RED = RGBColor(0xFF, 0x45, 0x45)
PURPLE = RGBColor(0xB0, 0x8A, 0xFF)
CYAN = RGBColor(0x00, 0xE5, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xD4, 0xE6, 0xF8)
DIM = RGBColor(0x78, 0x96, 0xB6)
VDIM = RGBColor(0x2E, 0x44, 0x5C)
UC_RED = RGBColor(0xE0, 0x00, 0x00)
UC_GLD = RGBColor(0xF5, 0xBF, 0x00)


def rect(l, t, w, h, fill):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def txt(l, t, w, h, text, pt, color, bold=False, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def add_image(path, l, t, w, h):
    return slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))


def section_bar(l, t, w, label, accent):
    rect(l, t, w, 1.10, CARD2)
    rect(l, t, 0.26, 1.10, accent)
    rect(l + 0.26, t, w - 0.26, 0.08, accent)
    txt(l + 0.52, t + 0.16, w - 0.70, 0.82, label, 48, accent, bold=True)
    return t + 1.10


def big_stat(l, t, w, h, value, label, accent):
    rect(l + 0.08, t + 0.08, w, h, VDIM)
    rect(l, t, w, h, CARD)
    rect(l, t, w, 0.14, accent)
    rect(l, t + h - 0.14, w, 0.14, accent)
    txt(l, t + 0.18, w, h * 0.52, value, 96, accent, bold=True, align=PP_ALIGN.CENTER)
    txt(l, t + h * 0.70, w, h * 0.26, label, 30, DIM, align=PP_ALIGN.CENTER)


def pill(l, t, w, h, text, bg, pt=32):
    rect(l + 0.07, t + 0.07, w, h, VDIM)
    rect(l, t, w, h, bg)
    txt(l, t + (h - pt / 72 * 1.4) / 2, w, h, text, pt, WHITE, bold=True, align=PP_ALIGN.CENTER)


def connector(l, t):
    txt(l, t, 0.55, 0.40, "▼", 28, VDIM, align=PP_ALIGN.CENTER)


def hbar(l, t, w, label, pct, color, max_pct=40):
    bw = w * 0.43
    txt(l, t, w * 0.43, 0.58, label, 30, OFFWHT)
    rect(l + w * 0.43, t + 0.16, bw, 0.24, VDIM)
    fill = bw * min(pct / max_pct, 1.0)
    if fill > 0.01:
        rect(l + w * 0.43, t + 0.16, fill, 0.24, color)
    txt(l + w * 0.43 + bw + 0.12, t + 0.04, 1.2, 0.52, f"{pct}%", 28, color, bold=True)


def quote_box(l, t, w, h, text, accent):
    rect(l, t, w, h, CARD)
    rect(l, t, 0.18, h, accent)
    txt(l + 0.35, t + 0.14, w - 0.50, h - 0.24, text, 32, accent, italic=True)


def extract_v4_assets() -> tuple[Path | None, Path | None]:
    if not KEYNOTE_PATH.exists():
        return None, None

    temp_dir = Path(tempfile.mkdtemp(prefix="draftvision_assets_"))
    headshot_path = temp_dir / "jared_headshot.jpeg"
    action_path = temp_dir / "football_action.png"

    with zipfile.ZipFile(KEYNOTE_PATH) as zf:
        for info in zf.infolist():
            name = info.filename
            if name.endswith("Facetune_16-02-2026-15-33-07-23.jpeg"):
                headshot_path.write_bytes(zf.read(info))
            elif "Screenshot 2026-04-06 at 10.49-54.png" in name:
                action_path.write_bytes(zf.read(info))

    return (
        headshot_path if headshot_path.exists() else None,
        action_path if action_path.exists() else None,
    )


HEADSHOT_PATH, ACTION_PATH = extract_v4_assets()

# ═════════════════════════════════════════════════════════════════════════════
#  BACKGROUND
# ═════════════════════════════════════════════════════════════════════════════
rect(0, 0, 42, 48, BG)
for i in range(0, 49, 4):
    rect(0, i, 42, 0.03, RGBColor(0x0A, 0x16, 0x28))

# ═════════════════════════════════════════════════════════════════════════════
#  HEADER   0 → 8.0"
# ═════════════════════════════════════════════════════════════════════════════
rect(0, 0, 42, 8.0, RGBColor(0x02, 0x05, 0x10))
rect(0, 7.72, 42, 0.10, BLUE)
rect(0, 7.82, 42, 0.06, GREEN)
rect(0, 7.88, 42, 0.12, BG)

txt(24.0, -1.0, 20.0, 10.0, "D", 640, RGBColor(0x06, 0x10, 0x22), bold=True)

if ACTION_PATH:
    add_image(ACTION_PATH, 28.0, 0.18, 13.6, 7.20)
    rect(28.0, 0.18, 3.20, 7.20, RGBColor(0x02, 0x05, 0x10))
    rect(30.8, 0.18, 2.0, 7.20, RGBColor(0x03, 0x07, 0x14))

rect(38.20, 0.28, 3.50, 1.55, UC_RED)
txt(38.20, 0.32, 3.50, 0.72, "UC", 60, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(38.20, 0.92, 3.50, 0.56, "CEAS", 30, WHITE, align=PP_ALIGN.CENTER)

txt(0.55, 0.05, 28.0, 3.30, "DraftVision", 178, BLUE, bold=True)
rect(0.55, 3.18, 21.0, 0.16, BLUE)
rect(0.55, 3.34, 14.0, 0.08, GREEN)

txt(0.55, 3.50, 27.0, 1.20, "Machine Learning–Powered NFL Prospect Success Prediction", 48, OFFWHT)
txt(0.55, 4.76, 27.0, 0.84, "Jared Krekeler  ·  B.S. Computer Science  ·  University of Cincinnati", 38, DIM)
txt(0.55, 5.62, 27.0, 0.75, "Advisor: Hrishi Vinayak Bhide  ·  Spring 2025  ·  CS Senior Design Expo", 32, VDIM)

# ═════════════════════════════════════════════════════════════════════════════
#  THREE COLUMNS  (8.2 → 28.0)
# ═════════════════════════════════════════════════════════════════════════════
GAP = 0.34
CW = 13.44
C1, C2, C3 = 0.40, 0.40 + CW + GAP, 0.40 + 2 * (CW + GAP)
R1T, R1H = 8.20, 19.80

for cx in [C1, C2, C3]:
    rect(cx, R1T, CW, R1H, PANEL)

# COL 1

y = section_bar(C1, R1T, CW, "INTRODUCTION", BLUE)
y += 0.38

txt(C1 + 0.38, y, CW - 0.76, 0.72, "The Problem", 40, BLUE, bold=True)
y += 0.82
txt(C1 + 0.38, y, CW - 0.76, 3.80, "NFL teams invest millions in draft picks, yet evaluating whether a college prospect will succeed at the pro level remains highly subjective. Traditional scouting relies on film study and instinct — with no standardized model to quantify success probability.", 32, OFFWHT)
y += 3.95

quote_box(C1 + 0.30, y, CW - 0.60, 1.80, '"Only ~31% of drafted players become long-term NFL starters."', GOLD)
y += 2.00

txt(C1 + 0.38, y, CW - 0.76, 0.72, "Project Objective", 40, BLUE, bold=True)
y += 0.82
txt(C1 + 0.38, y, CW - 0.76, 3.0, "DraftVision predicts the probability any college prospect will succeed in the NFL — powered by real ESPN live stats, combine athleticism, college tier, and draft capital signals.", 32, OFFWHT)
y += 3.18

txt(C1 + 0.38, y, CW - 0.76, 0.72, "Pipeline", 40, BLUE, bold=True)
y += 0.82

steps = [
    ("🏈  College Prospect", BLUE),
    ("📡  ESPN Live Stats", CYAN),
    ("⚙️   15-Feature Vector", PURPLE),
    ("🤖  XGBoost Classifier", BLUE),
    ("✅  Success Probability", GREEN),
]
ph, ch = 0.88, 0.38
for i, (label, bg) in enumerate(steps):
    pill(C1 + 0.42, y, CW - 0.84, ph, label, bg, 30)
    y += ph
    if i < len(steps) - 1:
        connector(C1 + CW / 2 - 0.28, y)
        y += ch

# COL 2

y = section_bar(C2, R1T, CW, "METHODS & DATA", PURPLE)
y += 0.38

txt(C2 + 0.38, y, CW - 0.76, 0.72, "Data Sources", 40, PURPLE, bold=True)
y += 0.82

for col, title, desc in [
    (CYAN, "ESPN CFB Roster API", "250 schools · 5,000+ athletes · live sync"),
    (BLUE, "ESPN Stats API", "Real season stats per player (ESPN IDs)"),
    (PURPLE, "Synthetic Training Labels", "4,000 samples — draft-round hit-rate labels"),
]:
    rect(C2 + 0.30, y, CW - 0.60, 1.18, CARD)
    rect(C2 + 0.30, y, 0.12, 1.18, col)
    rect(C2 + 0.30, y, CW - 0.60, 0.08, col)
    txt(C2 + 0.58, y + 0.13, CW - 1.0, 0.50, title, 30, col, bold=True)
    txt(C2 + 0.58, y + 0.65, CW - 1.0, 0.46, desc, 26, DIM)
    y += 1.30

y += 0.25
txt(C2 + 0.38, y, CW - 0.76, 0.72, "15 Input Features", 40, PURPLE, bold=True)
y += 0.82

for fname, fdesc, fc in [
    ("Draft Round", "Dominant — 34% of model gain", GOLD),
    ("Combine Speed Score", "Position-normalized athleticism 0–100", BLUE),
    ("College Tier", "Power 5 / Group of 5 / FCS", GREEN),
    ("Production Score", "Composite stat output 0–100", PURPLE),
    ("Volume Stats", "Pass yds/TDs, rush yds/TDs, games", CYAN),
    ("Position Flags ×5", "QB · RB · WR · TE · Other  (one-hot)", DIM),
]:
    rect(C2 + 0.30, y, CW - 0.60, 1.00, CARD)
    rect(C2 + 0.30, y, 0.10, 1.00, fc)
    txt(C2 + 0.54, y + 0.10, CW - 1.0, 0.46, fname, 30, fc, bold=True)
    txt(C2 + 0.54, y + 0.57, CW - 1.0, 0.38, fdesc, 25, DIM, italic=True)
    y += 1.10

y += 0.30
txt(C2 + 0.38, y, CW - 0.76, 0.72, "Why XGBoost?", 40, PURPLE, bold=True)
y += 0.82

for label, val, fc in [
    ("Logistic Regression", "72.0%", DIM),
    ("Random Forest", "74.2%", BLUE),
    ("XGBoost  ✓", "86.8%", GREEN),
]:
    bg_c = RGBColor(0x04, 0x24, 0x14) if fc == GREEN else CARD
    rect(C2 + 0.30, y, CW - 0.60, 0.88, bg_c)
    rect(C2 + 0.30, y, CW - 0.60, 0.09, fc)
    txt(C2 + 0.50, y + 0.14, 8.0, 0.50, label, 30, fc, bold=(fc == GREEN))
    bx = C2 + 0.30 + CW - 0.60 - 2.40
    rect(bx, y + 0.18, 2.28, 0.52, GREEN if fc == GREEN else VDIM)
    txt(bx, y + 0.22, 2.28, 0.50, val, 30, WHITE if fc == GREEN else DIM, bold=True, align=PP_ALIGN.CENTER)
    y += 0.98

# COL 3

y = section_bar(C3, R1T, CW, "RESULTS", GREEN)
y += 0.38

txt(C3 + 0.38, y, CW - 0.76, 0.72, "Model Performance", 40, GREEN, bold=True)
y += 0.82

bw = (CW - 0.60 - 0.28) / 3
for i, (val, lbl, vc) in enumerate([
    ("86.8%", "Accuracy", GREEN),
    ("0.928", "AUC", BLUE),
    ("0.85", "Precision", GOLD),
]):
    big_stat(C3 + 0.30 + i * (bw + 0.14), y, bw, 2.00, val, lbl, vc)
y += 2.22

txt(C3 + 0.38, y, CW - 0.76, 0.72, "Feature Importance (XGBoost Gain)", 40, GREEN, bold=True)
y += 0.82

for fname, imp, bc in [
    ("Draft Round", 34.3, GOLD),
    ("Combine Athleticism", 6.8, BLUE),
    ("Production Score", 5.5, GREEN),
    ("College Tier", 5.5, GREEN),
    ("Passing Yards", 5.4, DIM),
    ("Rushing Yards", 5.3, DIM),
    ("Rushing TDs", 5.0, DIM),
    ("Position: RB", 4.9, DIM),
]:
    hbar(C3 + 0.30, y, CW - 0.60, fname, imp, bc, max_pct=40)
    y += 0.60

y += 0.30
txt(C3 + 0.38, y, CW - 0.76, 0.72, "AUC Score", 40, GREEN, bold=True)
y += 0.82

rect(C3 + 0.30, y, CW - 0.60, 2.40, CARD)
rect(C3 + 0.30, y, CW - 0.60, 0.14, GREEN)
txt(C3 + 0.30, y + 0.20, CW - 0.60, 1.20, "0.928", 110, GREEN, bold=True, align=PP_ALIGN.CENTER)
txt(C3 + 0.30, y + 1.44, CW - 0.60, 0.55, "vs. random classifier = 0.50", 28, DIM, align=PP_ALIGN.CENTER, italic=True)
y += 2.60

txt(C3 + 0.38, y, CW - 0.76, 0.72, "Prediction Output", 40, GREEN, bold=True)
y += 0.82

for label, val, fc in [
    ("NFL Success", "68%  →  Round 1 pick", GREEN),
    ("Marginal", "26%  →  Round 4 pick", GOLD),
    ("Low Probability", "4%   →  Undrafted", RED),
]:
    rect(C3 + 0.30, y, CW - 0.60, 0.88, CARD)
    rect(C3 + 0.30, y, 0.12, 0.88, fc)
    txt(C3 + 0.52, y + 0.10, 4.0, 0.40, label, 26, DIM)
    txt(C3 + 0.52, y + 0.48, CW - 1.0, 0.38, val, 28, fc, bold=True)
    y += 0.98

# ═════════════════════════════════════════════════════════════════════════════
#  FULL-WIDTH: DRAFT ROUND DYNAMICS + SYSTEM APPLICATION
# ═════════════════════════════════════════════════════════════════════════════
SAT = R1T + R1H + 0.28
SAH = 7.70
rect(0.40, SAT, 41.20, SAH, PANEL)
section_bar(0.40, SAT, 41.20, "DRAFT ROUND DYNAMICS & SYSTEM APPLICATION", GREEN)

LW = 18.50
y2 = SAT + 1.22

txt(0.65, y2, LW - 0.35, 0.72, "Empirical NFL Success Rate by Draft Round", 38, GREEN, bold=True)
y2 += 0.85

for rnd, pct, col, label in [
    (1, 68, GREEN, "Round 1"),
    (2, 52, LBLUE, "Round 2"),
    (3, 38, BLUE, "Round 3"),
    (4, 26, PURPLE, "Round 4"),
    (5, 18, DIM, "Round 5"),
    (6, 12, DIM, "Round 6"),
    (7, 8, VDIM, "Round 7"),
    (8, 4, VDIM, "Undrafted"),
]:
    baw = LW - 0.50
    rh = 0.64
    rect(0.60, y2, baw, rh - 0.08, CARD)
    rect(0.60, y2, baw * pct / 100, rh - 0.08, col)
    txt(0.74, y2 + 0.10, 3.40, 0.46, label, 28, WHITE, bold=True)
    txt(0.60 + baw * pct / 100 + 0.16, y2 + 0.10, 1.80, 0.46, f"{pct}%", 28, col, bold=True)
    y2 += rh

RX = 0.40 + LW + 0.40
RW = 41.20 - LW - 0.80
y3 = SAT + 1.22

txt(RX, y3, RW, 0.72, "Web Platform  —  DraftVision", 38, GREEN, bold=True)
y3 += 0.85
txt(RX, y3, RW, 2.20, "Scouts, analysts, and fans enter any prospect's name and instantly receive an AI scouting report — success probability, live ESPN stats, combine profile, and model factor breakdown.", 32, OFFWHT)
y3 += 2.38

tech = [
    ("🐍 Flask + Python", GREEN),
    ("⚛️  React 18", BLUE),
    ("🤖 XGBoost", PURPLE),
    ("📡 ESPN API", CYAN),
    ("🗄️  PostgreSQL", GOLD),
    ("☁️  Railway", DIM),
]
pw = (RW - 0.30) / 3
ph3 = 0.86
for i, (label, col) in enumerate(tech):
    px = RX + (i % 3) * (pw + 0.15)
    py = y3 + (i // 3) * (ph3 + 0.18)
    rect(px, py, pw, ph3, CARD)
    rect(px, py, pw, 0.10, col)
    txt(px, py + 0.16, pw, ph3 - 0.20, label, 28, col, bold=True, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
#  ROW 3 — Challenges | Impact | Future Work
# ═════════════════════════════════════════════════════════════════════════════
R3T = SAT + SAH + 0.28
R3H = 8.50

for cx in [C1, C2, C3]:
    rect(cx, R3T, CW, R3H, PANEL)

for cx, label, accent, items in [
    (C1, "CHALLENGES", RED, [
        ("No Labeled Dataset", "Trained on 4,000 synthetic samples whose labels mirror empirical draft-round hit-rates from NFL history.", RED),
        ("Class Imbalance", "Only ~31% of players succeed. Addressed with stratified splits and XGBoost's built-in class weighting.", GOLD),
        ("Position Heterogeneity", "Solved with per-position production normalization so QB yards and WR yards are scored on the same 0–100 scale.", PURPLE),
    ]),
    (C2, "IMPACT", GREEN, [
        ("For NFL Franchises", "Data-driven second opinion — identify undervalued prospects in rounds 3–7 where evaluation variance is highest.", GREEN),
        ("For Fans & Analysts", "Accessible web tool explaining why a prospect is projected to succeed — no data science background needed.", CYAN),
        ("Self-Updating System", "Live ESPN roster sync auto-updates every college season with zero manual data maintenance.", BLUE),
    ]),
    (C3, "FUTURE WORK", GOLD, [
        ("Real Outcome Labels", "Replace synthetic data with historical college-to-NFL outcomes using Pro Football Reference Career AV scores.", GOLD),
        ("Position-Specific Models", "Train separate XGBoost models per position group to eliminate cross-position noise in the feature space.", GREEN),
        ("Player Tracking Data", "Add Next Gen Stats — separation, route sharpness, and speed — for a richer profile beyond combine proxy scores.", BLUE),
    ]),
]:
    y = section_bar(cx, R3T, CW, label, accent)
    y += 0.32
    for title, body, fc in items:
        rect(cx + 0.25, y, CW - 0.50, 0.09, fc)
        y += 0.18
        txt(cx + 0.25, y, CW - 0.50, 0.60, title, 32, fc, bold=True)
        y += 0.70
        txt(cx + 0.25, y, CW - 0.50, 1.80, body, 28, OFFWHT)
        y += 2.00

# ═════════════════════════════════════════════════════════════════════════════
#  FOOTER
# ═════════════════════════════════════════════════════════════════════════════
FT = R3T + R3H + 0.30
FH = 48.0 - FT

rect(0, FT, 42, FH, RGBColor(0x02, 0x05, 0x10))
rect(0, FT, 42, 0.12, BLUE)
rect(0, FT + 0.12, 42, 0.07, GREEN)

if HEADSHOT_PATH:
    add_image(HEADSHOT_PATH, 0.48, FT + 0.35, 2.60, 3.20)
else:
    rect(0.48, FT + 0.35, 2.60, 3.20, CARD)

txt(3.30, FT + 0.38, 14.5, 0.96, "Jared Krekeler", 60, WHITE, bold=True)
txt(3.30, FT + 1.38, 14.5, 0.64, "B.S. Computer Science", 36, DIM)
txt(3.30, FT + 2.04, 14.5, 0.62, "University of Cincinnati", 36, DIM)
txt(3.30, FT + 2.70, 14.5, 0.56, "krekeljm@mail.uc.edu", 32, BLUE)

rect(19.0, FT + 0.35, 0.08, 3.30, BLUE)

txt(19.70, FT + 0.38, 13.0, 0.55, "Faculty Advisor", 28, VDIM)
txt(19.70, FT + 0.96, 13.0, 0.88, "Hrishi Vinayak Bhide", 50, WHITE, bold=True)
txt(19.70, FT + 1.88, 13.0, 0.62, "University of Cincinnati", 36, DIM)
txt(19.70, FT + 2.52, 13.0, 0.56, "Department of Computer Science", 30, DIM)

rect(33.80, FT + 0.35, 0.08, 3.30, BLUE)

txt(34.50, FT + 0.38, 7.0, 0.55, "Built With", 28, VDIM)
for i, label in enumerate([
    "Python · XGBoost · Flask",
    "React 18 · Auth0 · PostgreSQL",
    "ESPN CFB API · Railway",
]):
    txt(34.50, FT + 1.00 + i * 0.68, 7.0, 0.62, label, 30, OFFWHT)

rect(0, 47.82, 42, 0.10, UC_RED)
rect(0, 47.92, 42, 0.08, UC_GLD)

txt(0, FT + FH - 0.54, 42, 0.50, "DraftVision  ·  Jared Krekeler  ·  University of Cincinnati  ·  CS Senior Design Expo  ·  Spring 2025", 26, VDIM, align=PP_ALIGN.CENTER)

prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
